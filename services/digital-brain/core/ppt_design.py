
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

class PPTDesignTheme:
    """Base class for PPT Themes"""
    def apply_cover(self, slide, title, subtitle=None):
        pass
    def apply_content(self, slide, title, page_index=1):
        pass
    def get_body_text_color(self):
        """Returns the RGBColor for body text."""
        return RGBColor(0, 0, 0) # Default Black

class CorporateLightTheme(PPTDesignTheme):
    """
    A Professional Light Theme featuring:
    - Clean White/Light Grey Background
    - Professional Blue Accents
    - Minimalist Layout
    """
    def __init__(self, prs):
        self.prs = prs
        self.width = prs.slide_width
        self.height = prs.slide_height
        
        # Palette
        self.c_bg = RGBColor(248, 250, 252)      # Slate 50 (Very light grey)
        self.c_card_bg = RGBColor(255, 255, 255) # White
        self.c_accent_1 = RGBColor(37, 99, 235)  # Blue 600
        self.c_accent_2 = RGBColor(14, 165, 233) # Sky 500
        self.c_text_main = RGBColor(15, 23, 42)  # Slate 900
        self.c_text_dim = RGBColor(71, 85, 105)  # Slate 600
        self.c_border = RGBColor(226, 232, 240)  # Slate 200

    def get_body_text_color(self):
        return self.c_text_dim

    def _draw_bg(self, slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.width, self.height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.c_bg
        bg.line.fill.background()
        
        # Accent shape (Triangle/Diagonal)
        shape = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(10), Inches(-2), Inches(6), Inches(6))
        shape.rotation = 45
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.c_accent_1
        shape.fill.transparency = 0.9
        shape.line.fill.background()

    def _move_to_front(self, slide, shape):
        slide.shapes._spTree.remove(shape._element)
        slide.shapes._spTree.append(shape._element)

    def apply_cover(self, slide, title_text, subtitle_text="Solution Proposal"):
        self._draw_bg(slide)
        
        # 1. Main Title Card (Left aligned, clean)
        # Vertical Bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.5), Inches(0.15), Inches(4.5))
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.c_accent_1
        bar.line.fill.background()
        
        # Title Placeholder Handling
        if slide.shapes.title:
            title_ph = slide.shapes.title
            self._move_to_front(slide, title_ph)
            
            title_ph.left = Inches(1.5)
            title_ph.top = Inches(1.5)
            title_ph.width = Inches(10)
            title_ph.height = Inches(2.5)
            
            if not title_ph.has_text_frame:
                title_ph.text_frame.text = title_text
            
            p = title_ph.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            p.font.name = "Arial Black"
            p.font.size = Pt(54)
            p.font.color.rgb = self.c_text_main
            
        # Subtitle
        try:
            sub_ph = None
            for s in slide.placeholders:
                if s.placeholder_format.idx == 1:
                    sub_ph = s
                    break
            if sub_ph:
                self._move_to_front(slide, sub_ph)
                sub_ph.left = Inches(1.5)
                sub_ph.top = Inches(4.2)
                sub_ph.width = Inches(10)
                sub_ph.height = Inches(1.5)
                
                if not sub_ph.has_text_frame:
                    sub_ph.text_frame.text = subtitle_text
                    
                p = sub_ph.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT
                p.font.name = "Arial"
                p.font.size = Pt(24)
                p.font.color.rgb = self.c_text_dim
        except: pass

    def apply_content(self, slide, title_text, page_index=1):
        self._draw_bg(slide)
        
        # 1. Header Line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(12.33), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = self.c_border
        line.line.fill.background()
        
        # 2. Accent Dot
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(0.45), Inches(0.15), Inches(0.15))
        dot.fill.solid()
        dot.fill.fore_color.rgb = self.c_accent_1
        dot.line.fill.background()

        # 3. Handle Title
        if slide.shapes.title:
            title_ph = slide.shapes.title
            self._move_to_front(slide, title_ph)
            title_ph.left = Inches(0.8)
            title_ph.top = Inches(0.3)
            title_ph.width = Inches(11)
            title_ph.height = Inches(0.8)
            
            p = title_ph.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            p.font.name = "Arial"
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = self.c_text_main

        # 4. Handle Body
        try:
            body_ph = None
            for s in slide.placeholders:
                if s.placeholder_format.idx == 1 or s.placeholder_format.type in [2, 7]:
                    body_ph = s
                    break
            
            if body_ph:
                self._move_to_front(slide, body_ph)
                # Slightly indented
                body_ph.left = Inches(0.8) 
                body_ph.top = Inches(1.5)
                body_ph.width = Inches(11.5)
                body_ph.height = Inches(5.5)
                if not body_ph.has_text_frame:
                    body_ph.text_frame.text = ""
        except: pass


class ModernTechTheme(PPTDesignTheme):
    """
    A High-End Tech Theme featuring:
    - Dark Navy Background (Slate 900)
    - Glassmorphism Cards (Semi-transparent backgrounds with borders)
    - Neon Accents (Cyan/Purple)
    - Geometric Decor (Hexagons/Dots)
    """
    def get_body_text_color(self):
        return RGBColor(240, 240, 240) # White
    
    def __init__(self, prs):
        self.prs = prs
        self.width = prs.slide_width
        self.height = prs.slide_height
        
        # Palette
        self.c_bg = RGBColor(15, 23, 42)      # Slate 900
        self.c_card_bg = RGBColor(30, 41, 59) # Slate 800
        self.c_accent_1 = RGBColor(6, 182, 212) # Cyan 500
        self.c_accent_2 = RGBColor(99, 102, 241) # Indigo 500
        self.c_text_main = RGBColor(248, 250, 252) # Slate 50
        self.c_text_dim = RGBColor(148, 163, 184) # Slate 400
        self.c_border = RGBColor(51, 65, 85) # Slate 700

    def _draw_bg(self, slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.width, self.height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.c_bg
        bg.line.fill.background()
        
        # Subtle texture (very large transparent circle)
        glow = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(-2), Inches(8), Inches(8))
        glow.fill.solid()
        glow.fill.fore_color.rgb = self.c_accent_2
        glow.fill.transparency = 0.95
        glow.line.fill.background()

    def _move_to_front(self, slide, shape):
        slide.shapes._spTree.remove(shape._element)
        slide.shapes._spTree.append(shape._element)

    def apply_cover(self, slide, title_text, subtitle_text="AI Solution Proposal"):
        self._draw_bg(slide)
        
        # 1. Glass Card for Title
        card_w = Inches(10)
        card_h = Inches(4)
        # Use simple integer division for EMU (integers)
        # prs.slide_width and Inches() are integers, but division returns float. Cast to int.
        card_x = int((self.width - card_w) / 2)
        card_y = int((self.height - card_h) / 2)
        
        # "Glass" Effect
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_x, card_y, card_w, card_h)
        card.fill.background() # Transparent
        card.line.color.rgb = self.c_accent_1
        card.line.width = Pt(1.5)
        
        # 2. Decor
        # Top-Right Tech Box
        dec1 = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, Inches(11.5), Inches(1), Inches(1.5), Inches(1.5))
        dec1.fill.solid()
        dec1.fill.fore_color.rgb = self.c_accent_1
        dec1.fill.transparency = 0.8
        
        # Bottom-Left Tech Dots
        for i in range(3):
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.5 + i*0.4), Inches(6), Inches(0.2), Inches(0.2))
            dot.fill.solid()
            dot.fill.fore_color.rgb = self.c_accent_2
            
        # 3. Handle Placeholders (Move to Front & Style)
        # Title
        if slide.shapes.title:
            title_ph = slide.shapes.title
            self._move_to_front(slide, title_ph)
            
            title_ph.left = card_x
            title_ph.top = card_y + Inches(1)
            title_ph.width = card_w
            title_ph.height = Inches(1.5)
            
            if not title_ph.has_text_frame:
                title_ph.text_frame.text = title_text
            
            p = title_ph.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.name = "Microsoft YaHei"
            p.font.size = Pt(44)
            p.font.bold = True
            p.font.color.rgb = self.c_text_main
            title_ph.text_frame.word_wrap = True # Wrap
        
        # Subtitle (usually idx 1 on Title layout)
        try:
            # Try to find subtitle placeholder (type 4 or idx 1)
            sub_ph = None
            for s in slide.placeholders:
                if s.placeholder_format.idx == 1:
                    sub_ph = s
                    break
            
            if sub_ph:
                self._move_to_front(slide, sub_ph)
                sub_ph.left = card_x
                sub_ph.top = card_y + Inches(2.5)
                sub_ph.width = card_w
                sub_ph.height = Inches(1)
                
                if not sub_ph.has_text_frame:
                    sub_ph.text_frame.text = subtitle_text
                    
                p_sub = sub_ph.text_frame.paragraphs[0]
                p_sub.alignment = PP_ALIGN.CENTER
                p_sub.font.name = "Arial"
                p_sub.font.size = Pt(20)
                p_sub.font.color.rgb = self.c_text_dim
        except Exception:
            pass # No subtitle placeholder found

    def apply_content(self, slide, title_text, page_index=1):
        # Override with Light Background for Content Slides (Modern Tech = Mixed Theme)
        # bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.width, self.height)
        # bg.fill.solid()
        # bg.fill.fore_color.rgb = RGBColor(248, 250, 252) # Slate 50
        # bg.line.fill.background()
        self._draw_bg(slide) # Back to Dark
        
        # 1. Header Area Decoration
        header_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.8), Inches(0.05), Inches(0.6))
        header_bar.fill.solid()
        header_bar.fill.fore_color.rgb = self.c_accent_1
        header_bar.line.fill.background()
        
        # 2. Main Body Card (White for Readability)
        body_x = Inches(0.5)
        body_y = Inches(1.8)
        body_w = self.width - Inches(1)
        body_h = self.height - Inches(2.5)
        
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, body_x, body_y, body_w, body_h)
        card.fill.background() # Transparent
        card.line.color.rgb = self.c_border
        card.line.width = Pt(1)
        
        # 3. Handle Placeholders
        # Title
        if slide.shapes.title:
            title_ph = slide.shapes.title
            self._move_to_front(slide, title_ph)
            
            title_ph.left = Inches(0.7)
            title_ph.top = Inches(0.5)
            title_ph.width = Inches(10)
            title_ph.height = Inches(1)
            
            p = title_ph.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            p.font.name = "Microsoft YaHei"
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = self.c_text_main
            
        # Body (Content)
        try:
            body_ph = None
            for s in slide.placeholders:
                # Match router logic: idx 1 OR type Body(2)/Object(7)
                if s.placeholder_format.idx == 1 or s.placeholder_format.type in [2, 7]:
                    body_ph = s
                    break
            
            if body_ph:
                self._move_to_front(slide, body_ph)
                
                body_ph.left = body_x + Inches(0.2)
                body_ph.top = body_y + Inches(0.2)
                body_ph.width = body_w - Inches(0.4)
                body_ph.height = body_h - Inches(0.4)
                
                if not body_ph.has_text_frame:
                    body_ph.text_frame.text = ""
        except Exception:
            pass
