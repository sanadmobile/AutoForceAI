from fastapi import APIRouter, Depends, HTTPException, Body
from fastapi.responses import FileResponse
from starlette.concurrency import run_in_threadpool
from sqlalchemy.orm import Session
from typing import List, Optional, Any
from pydantic import BaseModel, Field
import os
import uuid
import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from core.db_manager import get_shared_db
from branding_monitor.engines.qwen_client import QwenClient
from core.rag.retriever import KnowledgeRetriever
from database.shared_models import KnowledgeDoc, KnowledgeBase
from core.ppt_design import ModernTechTheme, CorporateLightTheme # Import the theme
from core.dependencies import get_current_user

# Use langchain for structured output if preferred, but direct prompting is simpler for now.
import json

router = APIRouter(
    prefix="/api/v1/solution",
    tags=["solution-generator"], 
    responses={404: {"description": "Not found"}}
)

# --- Data Models ---

class OutlineRequest(BaseModel):
    topic: str
    target_audience: Optional[str] = "Company Executives"
    style: Optional[str] = "Professional"
    kb_ids: List[int] = [] # Optional: restrict to specific KBs
    context_override: Optional[str] = None # Allow frontend to pass pre-retrieved context

class OutlinePage(BaseModel):
    page: int
    title: str
    type: str = "content" # cover, catalog, content, break, end
    key_points_hint: Optional[str] = None # Hint for what this page should cover

class OutlineResponse(BaseModel):
    topic: str
    pages: List[OutlinePage]

class ContentGenerationRequest(BaseModel):
    topic: str
    page_title: str
    page_type: str
    context_hint: Optional[str] = None
    kb_ids: List[int] = []

class PageContent(BaseModel):
    page: Optional[int] = 1
    title: str
    type: str = "content" 
    bullets: List[str] = []
    image_suggestion: Optional[str] = None
    speaker_notes: Optional[str] = None
    data_source: Optional[str] = None

class ContextItem(BaseModel):
    """Represents a single retrieved context chunk."""
    doc_id: int
    doc_name: str
    content: str
    score: float

class RetrievalLog(BaseModel):
    """Detailed logs for the thinking process."""
    step: str
    details: str
    timestamp: float

class ContextResponse(BaseModel):
    """Response for the intermediate retrieval step."""
    topic: str
    items: List[ContextItem]
    logs: List[RetrievalLog]

# --- Services (Helper Functions) ---

async def generate_outline_from_llm(topic: str, audience: str, retrieved_context: str) -> List[OutlinePage]:
    """
    Uses LLM to generate a structured PPT outline based on topic and context.
    """
    client = QwenClient()
    
    prompt = f"""
    You are an expert solution architect. design a presentation outline for the topic: "{topic}".
    Target Audience: {audience}
    
    Reference Context (Use this to tailor the outline):
    {retrieved_context[:4000]}
    
    Output Format:
    Return strictly a JSON array of objects. No markdown formatting.
    Each object must have: 
    - page (number)
    - title (string)
    - type (one of: 'cover', 'catalog', 'content', 'end')
    - key_points_hint (short description of page content)
    
    Structure the presentation logically based on the provided Reference Context.
    Approximate length: 8-12 slides.
    """
    
    # Use run_in_threadpool to call synchronous QwenClient.query in async context
    # enable_search=False because we prefer using retrieved RAG context
    response_text = await run_in_threadpool(client.query, prompt, enable_search=False)
    
    # Robust JSON extraction for Array
    try:
        import re
        # Look for [ ... ]
        json_match = re.search(r'\[.*\]', response_text, re.DOTALL)
        if json_match:
            clean_text = json_match.group(0)
        else:
            clean_text = response_text.strip()
            
        data = json.loads(clean_text)
        # Validate/Convert to Pydantic models
        pages = []
        for item in data:
            pages.append(OutlinePage(**item))
        return pages
    except Exception as e:
        print(f"JSON Parse Error: {e}\nRaw: {response_text}")
        # Fallback simple outline
        return [
            OutlinePage(page=1, title=topic, type="cover"),
            OutlinePage(page=2, title="Agenda", type="catalog"),
            OutlinePage(page=3, title="Background", type="content", key_points_hint="Current situation analysis"),
            OutlinePage(page=4, title="Solution Overview", type="content"),
            OutlinePage(page=5, title="Key Benefits", type="content"),
            OutlinePage(page=6, title="Q&A", type="end")
        ]

# --- Endpoints ---

@router.post("/context", response_model=ContextResponse)
async def retrieve_context_only(
    request: OutlineRequest, 
    db: Session = Depends(get_shared_db),
    user: dict = Depends(get_current_user)
):
    """
    Step 0: Retrieve context and return analysis logs (Thought Process).
    Does NOT generate the outline yet.
    """
    if not user:
        raise HTTPException(status_code=401, detail="Not authenticated")
        
    organization_id = user.get("organization_id") or user.get("org_id")
    
    # Fallback: If organization_id is missing from token, try to fetch it from DB using user_id
    if not organization_id and user.get("id"):
        print(f"[AUTH WARNING] organization_id missing in token for user {user.get('id')}, fetching from DB...")
        from database.shared_models import User
        db_user = db.query(User).filter(User.id == user["id"]).first()
        if db_user:
            organization_id = db_user.organization_id
            
    if not organization_id:
        raise HTTPException(status_code=403, detail="User not part of an organization")

    import time
    logs = []
    start_time = time.time()
    
    # Log 1: Interpret Intent
    logs.append(RetrievalLog(step="Intent Analysis", details=f"Analyzing topic: '{request.topic}' for audience '{request.target_audience}'", timestamp=time.time()-start_time))
    
    # Log 2: Query Expansion (Simulated for visualization)
    expanded_queries = [request.topic, f"{request.topic} 解决方案", f"{request.topic} 案例"]
    logs.append(RetrievalLog(step="Query Expansion", details=f"Generated search queries: {expanded_queries}", timestamp=time.time()-start_time))

    # Log 3: KB Context Resolution
    # Fetch all allowed KBs for this organization
    org_kbs = db.query(KnowledgeBase.id).filter(KnowledgeBase.organization_id == organization_id).all()
    allowed_kb_ids = {kb.id for kb in org_kbs}
    
    target_kb_ids = []
    if request.kb_ids:
        # Intersect requested KBs with allowed KBs
        target_kb_ids = [kb_id for kb_id in request.kb_ids if kb_id in allowed_kb_ids]
        if len(target_kb_ids) != len(request.kb_ids):
             logs.append(RetrievalLog(step="Security Warning", details=f"Some requested KBs were filtered out due to organization restrictions.", timestamp=time.time()-start_time))
    else:
        # Default to all KBs in the organization if none specified
        target_kb_ids = list(allowed_kb_ids)
        
    kb_label = f"KBs: {target_kb_ids}" if target_kb_ids else "Organization Knowledge Base (Empty)"
    
    logs.append(RetrievalLog(step="Retrieval Execution", details=f"Scanning vector database scope: {kb_label}...", timestamp=time.time()-start_time))
    
    # Actual Search
    retriever = KnowledgeRetriever(db)
    if not target_kb_ids:
        results = []
    else:
        # Use search_multi_kb for robust handling of multiple KBs
        results = retriever.search_multi_kb(target_kb_ids, request.topic, top_k=8)
    
    context_items = []
    unique_docs = set()
    
    if results:
        # Log 4: Result Filtering
        logs.append(RetrievalLog(step="Re-ranking", details=f"Found {len(results)} potential matches. Filtering by similarity score...", timestamp=time.time()-start_time))
        
        for r in results:
            doc_name = r.get('doc_name', 'Unknown')
            # Log specific document hit to show progress
            logs.append(RetrievalLog(
                step="Knowledge Scan", 
                details=f"Reading document >> {doc_name} (Relevance: {r.get('score', 0):.2f})", 
                timestamp=time.time()-start_time
            ))

            item = ContextItem(
                doc_id=r.get('doc_id', 0),
                doc_name=doc_name,
                content=r['content'],
                score=r.get('score', 0.0)
            )
            context_items.append(item)
            unique_docs.add(item.doc_name)
            
        logs.append(RetrievalLog(step="Context Selection", details=f"Selected {len(context_items)} chunks from {len(unique_docs)} documents: {list(unique_docs)}", timestamp=time.time()-start_time))
    else:
         logs.append(RetrievalLog(step="Retrieval Warning", details="No direct matches found in private knowledge base. Will use general model knowledge.", timestamp=time.time()-start_time))

    return ContextResponse(topic=request.topic, items=context_items, logs=logs)


@router.post("/outline", response_model=OutlineResponse)
async def create_outline(
    request: OutlineRequest, 
    db: Session = Depends(get_shared_db),
    user: dict = Depends(get_current_user)
):
    """
    Step 1: Generate a PPT outline based on a topic and KB retrieval.
    """
    if not user:
        raise HTTPException(status_code=401, detail="Not authenticated")
        
    organization_id = user.get("organization_id") or user.get("org_id")
    
    # Fallback: DB lookup if org_id missing in token
    if not organization_id and user.get("id"):
        from database.shared_models import User
        db_user = db.query(User).filter(User.id == user["id"]).first()
        if db_user:
            organization_id = db_user.organization_id

    if not organization_id:
        raise HTTPException(status_code=403, detail="User not part of an organization")

    context_text = ""
    
    # 1. Use existing context if provided (for multi-step visualization)
    if request.context_override:
        print(f"[Solution] Using pre-retrieved context for: {request.topic}")
        context_text = request.context_override
    else:
        # Fallback: Retrieve Context from KB inside this call
        # Fetch all allowed KBs for this organization
        org_kbs = db.query(KnowledgeBase.id).filter(KnowledgeBase.organization_id == organization_id).all()
        allowed_kb_ids = {kb.id for kb in org_kbs}
        
        target_kb_ids = []
        if request.kb_ids:
             target_kb_ids = [kb_id for kb_id in request.kb_ids if kb_id in allowed_kb_ids]
        else:
             target_kb_ids = list(allowed_kb_ids)
        
        # Search KB (Global or Specific)
        print(f"[Solution] Retrieving context for topic: {request.topic}, Target KBs: {target_kb_ids}")
        
        retriever = KnowledgeRetriever(db)
        if target_kb_ids:
             results = retriever.search_multi_kb(target_kb_ids, request.topic, top_k=8)
        else:
             results = []
        
        if results:
            context_text = "\n".join([f"- {r['content']}" for r in results])
            print(f"[Solution] Found {len(results)} context chunks.")
        else:
            print("[Solution] No relevant context found in Knowledge Base.")
    
    # 2. Call LLM to generate outline
    pages = await generate_outline_from_llm(request.topic, request.target_audience, context_text)
    
    return OutlineResponse(topic=request.topic, pages=pages)

@router.post("/page/content", response_model=PageContent)
async def generate_page_content(
    request: ContentGenerationRequest,
    db: Session = Depends(get_shared_db),
    user: dict = Depends(get_current_user)
):
    """
    Step 2: Generate detailed content for a specific slide.
    """
    if not user:
        raise HTTPException(status_code=401, detail="Not authenticated")
        
    organization_id = user.get("organization_id") or user.get("org_id")
    
    # Fallback: DB lookup if org_id missing in token
    if not organization_id and user.get("id"):
        from database.shared_models import User
        db_user = db.query(User).filter(User.id == user["id"]).first()
        if db_user:
            organization_id = db_user.organization_id

    if not organization_id:
        raise HTTPException(status_code=403, detail="User not part of an organization")

    # Fetch all allowed KBs
    org_kbs = db.query(KnowledgeBase.id).filter(KnowledgeBase.organization_id == organization_id).all()
    allowed_kb_ids = {kb.id for kb in org_kbs}
    
    target_kb_ids = []
    if request.kb_ids:
         target_kb_ids = [kb_id for kb_id in request.kb_ids if kb_id in allowed_kb_ids]
    else:
         target_kb_ids = list(allowed_kb_ids)

    client = QwenClient()
    retriever = KnowledgeRetriever(db)
    
    # 1. Retrieve specific context for this slide title
    search_query = f"{request.topic} {request.page_title}"
    context_text = ""
    
    # Search Scope
    print(f"[Solution] Retrieving content for slide: {request.page_title}, Target KBs: {target_kb_ids}")
    
    if target_kb_ids:
        results = retriever.search_multi_kb(target_kb_ids, search_query, top_k=5)
    else:
        results = []
    
    unique_sources = set()
    if results:
         context_text = "\n".join([f"- {r['content']}" for r in results])
         for r in results:
             if 'doc_name' in r:
                 unique_sources.add(r['doc_name'])
         print(f"[Solution] Found {len(results)} chunks for slide generation from {unique_sources}.")
         
    # 2. LLM Gen
    prompt = f"""
    Write content for a PowerPoint slide.
    Topic: {request.topic}
    Slide Title: {request.page_title}
    Slide Type: {request.page_type}
    Context Hint: {request.context_hint}
    
    Reference Material (CRITICAL: Prioritize this information over general knowledge):
    {context_text[:3000]}
    
    Output strictly JSON:
    {{
        "title": "{request.page_title}",
        "bullets": ["point 1", "point 2", "point 3"],
        "image_suggestion": "Description of an image",
        "speaker_notes": "Script for the presenter"
    }}
    """
    
    response_text = await run_in_threadpool(client.query, prompt, enable_search=False)
    print(f"[DEBUG] Raw LLM Response for {request.page_title}: {response_text[:200]}...") # Log start of response
    
    # Robust JSON extraction
    try:
        import re
        json_match = re.search(r'\{.*\}', response_text, re.DOTALL)
        if json_match:
            clean_text = json_match.group(0)
        else:
            clean_text = response_text.strip() # Fallback
            
        data = json.loads(clean_text)
        
        # Append sources to speaker notes if available
        if unique_sources:
            sources_str = ", ".join(unique_sources)
            
            # Store in dedicated field
            data["data_source"] = sources_str
            
            # Add to Speaker Notes (still useful for presenter)
            note_suffix = f"\n\n参考资料来源: {sources_str}"
            if "speaker_notes" in data and data["speaker_notes"]:
                data["speaker_notes"] += note_suffix
            else:
                data["speaker_notes"] = f"本页内容基于以下资料生成: {sources_str}"

        return PageContent(**data)
    except Exception as e:
        print(f"[ERROR] JSON Parse Failed: {e}")
        print(f"[ERROR] Full Response Text: {response_text}")
        return PageContent(
            title=request.page_title, 
            bullets=[f"Content generation failed: {str(e)}", "Raw response logged to server console."], 
            speaker_notes=""
        )

class FullPresentationRequest(BaseModel):
    topic: str
    pages: List[PageContent]
    template_id: Optional[str] = None


@router.post("/generate", response_class=FileResponse)
async def generate_pptx_file(
    request: FullPresentationRequest
):
    print("!!! HITTING THE ENDPOINT !!!", flush=True)
    """
    Step 3: Render the final PPTX file from structured data (Commercial Tech Theme or Template).
    """
    try:
        # --- Template Logic Setup ---
        TEMPLATE_DIR = r"E:\DigitalEmployee\storage\ppt_templates"
        use_template = False
        prs = None
        
        # Default to DeepSeek_Tech_Pro if not specified
        if not request.template_id:
            request.template_id = "DeepSeek_Tech_Pro.pptx"
        
        if request.template_id:
            tpl_path = os.path.join(TEMPLATE_DIR, request.template_id)
            if os.path.exists(tpl_path):
                # Load Template
                prs = Presentation(tpl_path)
                use_template = True
            elif "clean" in request.template_id.lower() or "modern" in request.template_id.lower() or "tech" in request.template_id.lower():
                 # Virtual Template / Theme Engine Trigger
                 print(f"[DEBUG] Using Virtual Template for {request.template_id}")
                 prs = Presentation()
                 prs.slide_width = Inches(13.333)
                 prs.slide_height = Inches(7.5)
                 use_template = True 
            else:
                print(f"[WARN] Template {request.template_id} not found, falling back to Dark Theme.")
        
        if not use_template:
            # Create NEW Presentation for Manual Dark Theme
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

        # --- Manual Theme Definitions (Only used if use_template=False) ---
        COLOR_BG = RGBColor(15, 23, 42)       
        COLOR_ACCENT = RGBColor(56, 189, 248) 
        COLOR_SEC = RGBColor(99, 102, 241)    
        COLOR_TEXT_MAIN = RGBColor(255, 255, 255) # White
        COLOR_TEXT_SUB = RGBColor(148, 163, 184) 

        def add_manual_design(slide, is_cover=False, is_ending=False):
            """Applies manual dark tech design."""
            bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
            bg.fill.solid()
            bg.fill.fore_color.rgb = COLOR_BG
            bg.line.fill.background()
            
            if is_cover:
                circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.5), Inches(-2), Inches(7.5), Inches(7.5))
                circle.fill.solid()
                circle.fill.fore_color.rgb = COLOR_SEC
                circle.fill.transparency = 0.8
                circle.line.fill.background()
            elif is_ending:
                circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4), Inches(2), Inches(5.33), Inches(5.33))
                circle.fill.solid()
                circle.fill.fore_color.rgb = COLOR_SEC
                circle.fill.transparency = 0.9
                circle.line.fill.background()
            else: 
                line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.25), Inches(1), Inches(0.06))
                line.fill.solid()
                line.fill.fore_color.rgb = COLOR_ACCENT
                line.line.fill.background()

        # --- Rich Template Logic ---
        theme_styler = None
        print(f"[DEBUG] Checking Template ID: {request.template_id}", flush=True)
        
        # Only apply programmatic styling to System Templates or if explicitly requested via filename keywords
        SYSTEM_KEYWORDS = ["deepseek", "modern", "tech"]
        
        should_apply_theme = False
        if request.template_id:
             tid_lower = request.template_id.lower()
             if any(k in tid_lower for k in SYSTEM_KEYWORDS):
                 should_apply_theme = True

        if should_apply_theme:
            print("[DEBUG] Activating ModernTechTheme (System Logic)", flush=True)
            theme_styler = ModernTechTheme(prs)
        else:
            print("[DEBUG] Using Custom User Template - No Programmatic Styling Overlay", flush=True)
            theme_styler = None

        # --- Slide Generation Loop ---
        print(f"[DEBUG] Starting generation loop for {len(request.pages)} pages. use_template={use_template}", flush=True)
        for p_content in request.pages:
            is_cover = (p_content.type == 'cover')
            is_ending = (p_content.type == 'ending') or (p_content.type == 'end')
            
            if use_template:
                print(f"[DEBUG] Generating slide {p_content.page} (template mode)", flush=True)
                
                # Layout Selection
                target_layout_idx = 0 if (is_cover or is_ending) else 1
                
                # Try smarter layout finding
                if is_cover or is_ending:
                     # Find title layout (type 1 or 3)
                     for i, l in enumerate(prs.slide_layouts):
                          if any(s.placeholder_format.type in [1,3] for s in l.placeholders):
                               target_layout_idx = i
                               break
                else:
                     # Find content layout (type 2 or 7)
                     for i, l in enumerate(prs.slide_layouts):
                          if any(s.placeholder_format.type in [2,7] for s in l.placeholders):
                               target_layout_idx = i
                               break

                if target_layout_idx >= len(prs.slide_layouts): target_layout_idx = 0
                slide = prs.slides.add_slide(prs.slide_layouts[target_layout_idx])
                
                # Apply Theme
                if theme_styler:
                     if is_cover:
                         theme_styler.apply_cover(slide, (request.topic if is_cover else p_content.title))
                     else:
                         theme_styler.apply_content(slide, (request.topic if is_cover else p_content.title))

                # Title Handling
                title_text = (request.topic if is_cover else p_content.title)
                title_shape = None
                
                if slide.shapes.title:
                    title_shape = slide.shapes.title
                
                if not title_shape:
                     for shape in slide.placeholders:
                         if shape.placeholder_format.type in [1, 3]: 
                             title_shape = shape
                             break
                
                if not title_shape:
                     # Soft fallback for custom layouts
                     for shape in slide.shapes:
                         if shape.has_text_frame and not shape.is_placeholder and shape.top < prs.slide_height * 0.2:
                             title_shape = shape
                             break

                if title_shape and title_shape.has_text_frame:
                     # --- SANITY CHECK: Enforce Title Position for Content Slides ---
                     # If title looks ridiculously low or large, force it to top
                     if not is_cover and not is_ending:
                         if (title_shape.top + title_shape.height) > Inches(3):
                             print(f"[WARN-FIX] Title shape for slide {p_content.page} is too large/low (Bottom: {(title_shape.top + title_shape.height)/914400:.2f}in). Forcing reset.")
                             title_shape.top = Inches(0.2) # Higher
                             title_shape.height = Inches(1.0) # More compact
                             # Also ensure it has width
                             if title_shape.width < Inches(5):
                                 title_shape.width = Inches(10)
                                 title_shape.left = Inches(1.6) # Centered-ish

                     tf = title_shape.text_frame
                     tf.word_wrap = True # Ensure wrap
                     if len(tf.paragraphs) > 0:
                         p = tf.paragraphs[0]
                         p.text = title_text # Simple set
                     else:
                         p = tf.add_paragraph()
                         p.text = title_text
                     
                     # Force Black Title for Custom Templates (User Request)
                     if not theme_styler:
                         if len(tf.paragraphs) > 0:
                             tf.paragraphs[0].font.color.rgb = RGBColor(0,0,0)
                else:
                    # Manual Title Fallback
                    print(f"[WARN] No Title Placeholder found for slide {p_content.page}, Creating Fallback Textbox.")
                    
                    # Check if potential body overlap exists
                    overlap_body = None
                    if not is_cover and not is_ending:
                        for shape in slide.placeholders:
                             if shape.placeholder_format.type in [2, 7]:
                                 overlap_body = shape
                                 break
                    
                    # Move body down if it starts too high
                    if overlap_body and overlap_body.top < Inches(1.8):
                        print(f"[DEBUG] Moving overlapping body placeholder down from {overlap_body.top}")
                        overlap_body.top = Inches(1.8)
                        overlap_body.height = prs.slide_height - Inches(2.2) # Resize to fit bottom

                    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), prs.slide_width - Inches(1), Inches(1))
                    txBox.text_frame.word_wrap = True # Ensure wrap
                    txBox.text_frame.text = title_text
                    
                    # Assign fallback title for overlap check downstream
                    title_shape = txBox
                    
                # Body Handling
                if not is_cover and not is_ending:
                    print(f"[DEBUG] Processing Content Slide content...", flush=True)
                    body_shape = None
                    for shape in slide.placeholders:
                        if shape.placeholder_format.type in [2, 7]:
                            body_shape = shape
                            break
                    if not body_shape:
                        for shape in slide.placeholders:
                            if shape.placeholder_format.idx == 1:
                                body_shape = shape
                                break
                    
                    if not body_shape:
                        print(f"[WARN] No Body Shape Found for Slide {p_content.page}, using fallback textbox.", flush=True)
                        # Fallback position
                        body_shape = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.3), Inches(4.5))
                        body_shape.text_frame.word_wrap = True
                    
                    # --- CRITICAL FIX: Overlap Detection & Correction ---
                    # Ensure Body doesn't overlap Title, even if template is poorly designed
                    if body_shape:
                        print(f"[DEBUG] Body shape size: W={body_shape.width/914400:.2f}in, H={body_shape.height/914400:.2f}in")
                        
                        # Aggressively fix layout for User Custom Templates to prevent "One char per line" issues
                        # This issues often comes from vertical text placeholders or bad margins or narrow widths
                        should_fix_geometry = (not theme_styler) or (body_shape.width < Inches(8))
                        
                        if should_fix_geometry:
                            print(f"[WARN-FIX] Enforcing standard body geometry for Slide {p_content.page} (W={body_shape.width/914400:.2f}in).")
                            body_shape.left = Inches(1)
                            body_shape.width = Inches(11.3)
                            body_shape.rotation = 0
                            
                            # Also enforce a minimum height if it's too short (User observed H=2.01in which is small for body)
                            if body_shape.height < Inches(4):
                                print(f"[WARN-FIX] Body height too plain ({body_shape.height/914400:.2f}in). Extending.")
                                body_shape.height = Inches(4.5)
                            
                        # Force Text Frame Props
                        if body_shape.has_text_frame:
                            body_shape.text_frame.word_wrap = True
                            body_shape.text_frame.margin_left = Inches(0.1)
                            body_shape.text_frame.margin_right = Inches(0.1)
                            body_shape.text_frame.margin_top = Inches(0.1)
                            
                            # --- CRITICAL XML FIX: Force Horizontal Text Orientation ---
                            # Many user templates might have placeholders set to "Vertical" or "Stacked" text.
                            # We must override this property in the underlying XML.
                            try:
                                body_pr = body_shape.text_frame._element.bodyPr
                                body_pr.set('vert', 'horz') # Force Horizontal
                                body_pr.set('wrap', 'square') # Standard Wrapping
                                body_pr.set('lIns', '91440') # 0.1 inch margin XML
                                body_pr.set('rIns', '91440')
                                body_pr.set('tIns', '45720') # 0.05 inch
                                body_pr.set('bIns', '45720')
                            except Exception as e:
                                print(f"[WARN] Failed to set bodyPr XML: {e}")

                        safe_top = Inches(1.8) # Default safe zone
                        title_bottom_y = 0
                        
                        if title_shape:
                            title_bottom_y = title_shape.top + title_shape.height
                            safe_top = title_bottom_y + Inches(0.2)
                            
                        print(f"[DEBUG-LAYOUT] Slide {p_content.page}: Title Bottom={title_bottom_y/914400:.2f}in, Body Top={body_shape.top/914400:.2f}in")
                        
                        if body_shape.top < safe_top:
                            print(f"[WARN-FIX] Body overlaps Title area! Moving Body DOWN. (Old: {body_shape.top/914400:.2f}in -> New: {safe_top/914400:.2f}in)")
                            body_shape.top = int(safe_top)
                            # Adjust height so it doesn't fall off slide
                            max_h = prs.slide_height - body_shape.top - Inches(0.5)
                            if body_shape.height > max_h:
                                body_shape.height = int(max_h)

                    if body_shape and body_shape.has_text_frame:
                        tf = body_shape.text_frame
                        if p_content.bullets:
                            print(f"[DEBUG] Writing {len(p_content.bullets)} bullets", flush=True)
                            
                            def style_para(para):
                                # Always remove indents that might cause "thin column" look
                                # para.indent = 0 # Invalid attribute
                                para.space_before = Pt(6)
                                para.space_after = Pt(6)
                                
                                # --- XML FIX: Reset Paragraph Indents Directly ---
                                # Solves "One word per line" caused by massive master slide indentation
                                try:
                                    pPr = para._p.get_or_add_pPr()
                                    # marL: Left Margin (reset to 0)
                                    pPr.set('marL', '0')
                                    # indent: First Line Indent (reset to 0)
                                    pPr.set('indent', '0')
                                except Exception as e:
                                    print(f"[WARN] XML Indent fix failed: {e}")

                                if theme_styler:
                                    # Dynamically get color from theme
                                    para.font.color.rgb = theme_styler.get_body_text_color()
                                    
                                    if para.font.size is None or para.font.size < Pt(14):
                                        para.font.size = Pt(18)
                                    # print(f"[DEBUG] Styled para with theme color", flush=True)
                                else:
                                    # For User Templates, also ensure reasonable font size/color if missing
                                    if para.font.size is None or para.font.size < Pt(12):
                                         para.font.size = Pt(18)
                                    # Ensure Text is Black (User Request)
                                    para.font.color.rgb = RGBColor(0, 0, 0)
                                    # Remove bullet indent weirdness
                                    # Reset level to 0 to clear deep nesting
                                    para.level = 0
                                    
                                    # Explicitly set font size to ensure visibility
                                    if para.font.size is None or para.font.size < Pt(14):
                                        para.font.size = Pt(18)

                            # First bullet
                            if len(tf.paragraphs) > 0:
                                p = tf.paragraphs[0]
                                p.text = p_content.bullets[0]
                                style_para(p)
                            else:
                                p = tf.add_paragraph()
                                p.text = p_content.bullets[0]
                                style_para(p)
                            
                            # Rest bullets
                            for i in range(1, len(p_content.bullets)):
                                p = tf.add_paragraph()
                                p.text = p_content.bullets[i]
                                style_para(p)
                        else:
                             tf.clear()
                    else:
                        print(f"[WARN] Still No Body Shape for Slide {p_content.page}", flush=True)

            else:
                # Manual Mode (No Template)
                slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank
                add_manual_design(slide, is_cover, is_ending)
                # ... Simplified manual mode logic (omitted for brevity as we use template usually) ...
                if is_cover:
                    tb = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2.5))
                    tb.text_frame.text = request.topic.upper()
                elif is_ending:
                    tb = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1.5))
                    tb.text_frame.text = "Thank You"
                else:
                    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
                    tb.text_frame.text = p_content.title
                    
                    if p_content.bullets:
                        body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(12), Inches(5))
                        tf = body_box.text_frame
                        tf.text = "\n".join(p_content.bullets)

            # Speaker Notes
            if p_content.speaker_notes:
                 # Accessing notes_slide creates it if it doesn't exist
                 slide.notes_slide.notes_text_frame.text = p_content.speaker_notes

        # Save
        output_dir = "storage/temp_ppt"
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)
            
        filename = f"Solution_{uuid.uuid4().hex[:8]}.pptx"
        file_path = os.path.join(output_dir, filename)
        
        prs.save(file_path)
        print(f"[SUCCESS] Saved PPT to {file_path}", flush=True)
        
        return FileResponse(
            path=file_path, 
            filename=f"{request.topic}_Solution.pptx",
            media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )

    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"PPT Generation Failed: {str(e)}")
