from fastapi import APIRouter, HTTPException, Depends, Body
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
import io
import os
import markdown
from bs4 import BeautifulSoup, NavigableString

router = APIRouter(prefix="/api/v1/export", tags=["Export"])

class ExportRequest(BaseModel):
    title: str = "Chat Export"
    content: str

def add_html_to_doc(doc, soup):
    """
    Parses BeautifulSoup object and adds content to python-docx Document
    """
    from docx.shared import Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    
    # Helper to process inline tags like <b>, <i>, <code> within a paragraph
    def process_inline(paragraph, element):
        if isinstance(element, NavigableString):
            run = paragraph.add_run(str(element))
            return run
        
        for child in element.children:
            if isinstance(child, NavigableString):
                run = paragraph.add_run(str(child))
                if element.name == 'strong' or element.name == 'b':
                    run.bold = True
                if element.name == 'em' or element.name == 'i':
                    run.italic = True
                if element.name == 'code':
                    run.font.name = 'Courier New'
                    run._element.rPr.rFonts.set(qn('w:eastAsia'), 'Courier New')
                    run.font.size = Pt(10)
            else:
                process_inline(paragraph, child)

    # Main loop for block elements
    for element in soup.children:
        if isinstance(element, NavigableString):
            text = str(element).strip()
            if text:
                doc.add_paragraph(text)
            continue
            
        if element.name in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
            level = int(element.name[1])
            if level == 1:
                doc.add_heading(element.get_text(), level=1)
            elif level == 2:
                doc.add_heading(element.get_text(), level=2)
            else:
                doc.add_heading(element.get_text(), level=level)
                
        elif element.name == 'p':
            p = doc.add_paragraph()
            # Process inline elements (b, i, code, etc.)
            for child in element.children:
                if isinstance(child, NavigableString):
                    p.add_run(str(child))
                elif child.name in ['strong', 'b']:
                    run = p.add_run(child.get_text())
                    run.bold = True
                elif child.name in ['em', 'i']:
                    run = p.add_run(child.get_text())
                    run.italic = True
                elif child.name == 'code':
                    run = p.add_run(child.get_text())
                    run.font.name = 'Courier New'
                    from docx.shared import Pt
                    run.font.size = Pt(10)
                else:
                    p.add_run(child.get_text())
            
        elif element.name == 'ul':
            for li in element.find_all('li', recursive=False):
                doc.add_paragraph(li.get_text(), style='List Bullet')
                
        elif element.name == 'ol':
            for li in element.find_all('li', recursive=False):
                doc.add_paragraph(li.get_text(), style='List Number')
                
        elif element.name == 'pre':
            # Code block
            code_text = element.get_text()
            p = doc.add_paragraph(code_text)
            from docx.shared import Pt, RGBColor
            for run in p.runs:
                run.font.name = 'Courier New'
                run.font.size = Pt(9)
                run.font.color.rgb = RGBColor(50, 50, 50)
            # Add a border or background if possible, but python-docx native styling is limited for paragraphs
            
        elif element.name == 'table':
            rows = element.find_all('tr')
            if not rows:
                continue
            
            # Count max columns
            max_cols = 0
            for row in rows:
                cols = row.find_all(['th', 'td'])
                if len(cols) > max_cols:
                    max_cols = len(cols)
            
            if max_cols > 0:
                table = doc.add_table(rows=len(rows), cols=max_cols)
                table.style = 'Table Grid'
                
                for i, row in enumerate(rows):
                    cols = row.find_all(['th', 'td'])
                    for j, col in enumerate(cols):
                        if j < max_cols:
                            cell = table.cell(i, j)
                            cell.text = col.get_text().strip()
                            # Helper to bold header row
                            if row.find('th'):
                                for paragraph in cell.paragraphs:
                                    for run in paragraph.runs:
                                        run.bold = True
        
        elif element.name == 'hr':
            doc.add_paragraph('_' * 20)

@router.post("/word")
def export_word(request: ExportRequest):
    try:
        from docx import Document
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Missing dependency python-docx: {e}")
    
    doc = Document()
    doc.add_heading(request.title, 0)
    
    # Convert Markdown to HTML
    try:
        html = markdown.markdown(
            request.content, 
            extensions=['tables', 'fenced_code', 'nl2br', 'sane_lists']
        )
        
        # Parse HTML
        soup = BeautifulSoup(html, 'html.parser')
        
        # Add to Doc
        add_html_to_doc(doc, soup)
    except Exception as e:
        print(f"Error parsing markown to docx: {e}")
        # Fallback
        doc.add_paragraph(request.content)
    
    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    
    return StreamingResponse(
        buffer, 
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers={"Content-Disposition": f"attachment; filename={request.title}.docx"}
    )


@router.post("/pdf")
def export_pdf(request: ExportRequest):
    try:
        from fpdf import FPDF
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Missing dependency fpdf: {e}")
    pdf = FPDF()
    pdf.add_page()
    
    # Try to load a Chinese font for Windows environment
    font_path = "C:\\Windows\\Fonts\\simhei.ttf"
    has_font = False
    
    if os.path.exists(font_path):
        try:
            pdf.add_font('SimHei', '', font_path, uni=True)
            pdf.set_font('SimHei', '', 12)
            has_font = True
        except Exception as e:
            print(f"Font loading error: {e}")
            
    if not has_font:
        # Fallback to Arial (Chinese will be garbled)
        pdf.set_font("Arial", size=12)
        
    # Sanitize content for PDF if needed
    txt = request.content
    
    # FPDF multi_cell for text wrapping
    pdf.multi_cell(0, 10, txt=txt)
    
    buffer = io.BytesIO()
    # Output to buffer (fpdf2 style)
    try:
        pdf.output(buffer)
    except TypeError:
        # FPDF 1.7.x style (returns string)
        res = pdf.output(dest='S')
        if isinstance(res, str):
            res = res.encode('latin-1')
        buffer.write(res)

    buffer.seek(0)
    return StreamingResponse(
        buffer,
        media_type="application/pdf",
        headers={"Content-Disposition": f"attachment; filename={request.title}.pdf"}
    )

@router.post("/ppt")
def export_ppt(request: ExportRequest):
    try:
        from pptx import Presentation
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Missing dependency python-pptx: {e}")
    prs = Presentation()
    # Layout 1 is usually Title + Content
    slide = prs.slides.add_slide(prs.slide_layouts[1]) 
    
    title = slide.shapes.title
    title.text = request.title
    
    # Placeholders[1] is usually the content body
    if len(slide.placeholders) > 1:
        content = slide.placeholders[1]
        content.text = request.content
    else:
        # Fallback if layout is different
        txBox = slide.shapes.add_textbox(left=0, top=100, width=500, height=500)
        tf = txBox.text_frame
        tf.text = request.content
    
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    
    return StreamingResponse(
        buffer,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename={request.title}.pptx"}
    )

@router.post("/excel")
def export_excel(request: ExportRequest):
    try:
        import pandas as pd
    except Exception as e:
         raise HTTPException(status_code=500, detail=f"Missing dependency pandas: {e}")
            
    buffer = io.BytesIO()
    
    # Simple markdown table parser strategy
    # 1. We look for blocks of lines starting and ending with |
    # 2. We extract them as DataFrames
    # 3. We put everything into an Excel file
    
    lines = request.content.split('\n')
    data_frames = []
    current_table_lines = []
    
    def flush_table(table_lines, dfs):
        if len(table_lines) < 2: return 
        
        # Filter out separator lines (e.g. |---|---|)
        # Regex-like check: if line contains only | - : and whitespace
        def is_separator(l):
            chars = set(l.replace('|', '').replace('-', '').replace(':', '').strip())
            return not chars
            
        data_lines = [l for l in table_lines if not is_separator(l)]
        
        if len(data_lines) < 1: return

        try:
            # Manual parse to handle outer pipes elegantly
            rows = []
            for l in data_lines:
                # Split by |
                # Typical markdown table: | Col1 | Col2 |
                # split('|') -> ['', ' Col1 ', ' Col2 ', '']
                parts = l.strip().split('|')
                # If starts/ends with pipe, the first/last elements are empty strings, remove them
                if l.strip().startswith('|'): parts.pop(0)
                if l.strip().endswith('|') and len(parts) > 0: parts.pop(-1)
                
                rows.append([p.strip() for p in parts])
            
            if not rows: return
            
            # Assume first row is header
            header = rows[0]
            data = rows[1:] if len(rows) > 1 else []
            
            # Pad data rows if length mismatch
            for r in data:
                while len(r) < len(header): r.append("")
                while len(r) > len(header): r.pop()
                    
            df = pd.DataFrame(data, columns=header)
            dfs.append(df)
        except Exception as e:
            print(f"Table parse warning: {e}")

    for line in lines:
        stripped = line.strip()
        # Basic detection of table line
        if stripped.startswith('|'):
            current_table_lines.append(stripped)
        else:
            if current_table_lines:
                flush_table(current_table_lines, data_frames)
                current_table_lines = []
    
    # Flush any remaining
    if current_table_lines:
        flush_table(current_table_lines, data_frames)

    # Write to Excel
    # If no tables found, put the raw text in cell A1
    with pd.ExcelWriter(buffer, engine='openpyxl') as writer:
        if data_frames:
            for i, df in enumerate(data_frames):
                sheet_name = f"Table_{i+1}"
                try:
                    df.to_excel(writer, sheet_name=sheet_name, index=False)
                except Exception as e:
                    # In case of invalid sheet name or other error
                    df.to_excel(writer, sheet_name=f"Sheet{i+1}", index=False)
            
            # Also add a "Full Content" sheet just in case
            pd.DataFrame({'Full Text': [request.content]}).to_excel(writer, sheet_name="Source_Text", index=False)
        else:
            pd.DataFrame({'Content': [request.content]}).to_excel(writer, sheet_name="Chart_Export", index=False)

    buffer.seek(0)
    
    return StreamingResponse(
        buffer,
         media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": f"attachment; filename={request.title}.xlsx"}
    )
