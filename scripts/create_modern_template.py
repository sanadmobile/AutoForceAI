
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_THEME_COLOR

def create_modern_template():
    prs = Presentation()
    # 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Colors ---
    # DEEP_BLUE = RGBColor(10, 25, 47)
    # NEON_TEAL = RGBColor(100, 255, 218)
    # SLATE_WHITE = RGBColor(230, 241, 255)
    
    # --- 1. Title Layout (Index 0) ---
    # We edit the slide master layouts directly or just create a blank one and add placeholders?
    # python-pptx creates a default set. We will modify the first two layouts of the slide master.
    
    # Get the master
    master = prs.slide_master
    
    # --- Customize Title Layout (master.slide_layouts[0]) ---
    title_layout = master.slide_layouts[0]
    
    # 1. Background (Dark)
    bg = title_layout.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 23, 42) # Dark Slate

    # 2. Add "Tech" Decorations - REMOVED due to python-pptx limitation on Layouts
    # Hexagon
    # shapes = title_layout.shapes
    # hex1 = shapes.add_auto_shape(MSO_SHAPE.HEXAGON, Inches(10), Inches(-1), Inches(4), Inches(4)) ...
    
    # 3. Style Title Placeholder
    # We need to find the existing placeholders or ensure they exist
    # Standard Title Layout has Title (idx 0) and Subtitle (idx 1)
    
    if title_layout.placeholders:
        title_ph = title_layout.placeholders[0] # Usually Title
        title_ph.left = Inches(1)
        title_ph.top = Inches(2.5)
        title_ph.width = Inches(10)
        title_ph.height = Inches(2)
        
        tf = title_ph.text_frame
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(54)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.bold = True
        
        # Subtitle
        if len(title_layout.placeholders) > 1:
            sub = title_layout.placeholders[1]
            sub.left = Inches(1)
            sub.top = Inches(4.6)
            sub.width = Inches(10)
            sub.height = Inches(1)
            
            sub.text_frame.paragraphs[0].font.color.rgb = RGBColor(148, 163, 184) # Slate 400
            sub.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            sub.text_frame.paragraphs[0].font.name = "Microsoft YaHei"


    # --- 2. Content Layout (Visual/Card Style) (Index 1) ---
    # We want a layout that looks like a "Card" or has a "Sidebar"
    content_layout = master.slide_layouts[1]
    
    # Background: Light Gray
    c_bg = content_layout.background
    c_bg.fill.solid()
    c_bg.fill.fore_color.rgb = RGBColor(248, 250, 252) # Slate 50
    
    # Add a Visual Header Bar (Gradient-like using shapes) - REMOVED
    # header_bar = content_layout.shapes.add_auto_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.4), Inches(0.15))
    # ...

    # Style Title Placeholder
    # Find existing title
    c_title = None
    c_body = None
    for s in content_layout.placeholders:
        if s.placeholder_format.idx == 0: c_title = s
        if s.placeholder_format.idx == 1: c_body = s
        
    if c_title:
        c_title.left = Inches(0.8)
        c_title.top = Inches(0.5)
        c_title.width = Inches(11)
        c_title.height = Inches(1)
        c_title.text_frame.paragraphs[0].font.name = "Microsoft YaHei"
        c_title.text_frame.paragraphs[0].font.size = Pt(32)
        c_title.text_frame.paragraphs[0].font.color.rgb = RGBColor(15, 23, 42) # Slate 900
        c_title.text_frame.paragraphs[0].font.bold = True
        c_title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    # Style Body Placeholder -> Make it look like a "Card"
    # Note: Placeholder shapes can have fills!
    if c_body:
        c_body.left = Inches(0.8)
        c_body.top = Inches(1.8)
        c_body.width = Inches(11.7)
        c_body.height = Inches(5)
        
        # Visual: White Card Background with gentle Border
        fill = c_body.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        line = c_body.line
        line.color.rgb = RGBColor(226, 232, 240) # Slate 200
        line.width = Pt(1)
        
        # Add a shadow effect? python-pptx limited here, but fill/line works.
        
        # Important: Text Padding
        c_body.text_frame.margin_left = Inches(0.5)
        c_body.text_frame.margin_right = Inches(0.5)
        c_body.text_frame.margin_top = Inches(0.5)
        c_body.text_frame.margin_bottom = Inches(0.5)

        # Draw a "Connection Visual" on the left of the text
        # To simulate "Relationship", we add a vertical line with list dots decoration ON TOP of the placeholder area? 
        # No, shapes added to layout appear behind placeholders usually. 
        # Let's add a decorative sidebar to the layout
        
        # Visual Node Connector Sidebar - REMOVED
        # sidebar = content_layout.shapes.add_auto_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(2.3), Inches(0.05), Inches(4))
        # ...
        
        # Add some dots - REMOVED
        # for y in [2.5, 3.5, 4.5, 5.5]:
        #    dot = content_layout.shapes.add_auto_shape(MSO_SHAPE.OVAL, Inches(0.85), Inches(y), Inches(0.35), Inches(0.35))
        #    ...


    # --- 3. Save ---
    output_path = r"E:\DigitalEmployee\storage\ppt_templates\Modern_Tech_Workflow.pptx"
    prs.save(output_path)
    print(f"Created template at: {output_path}")

if __name__ == "__main__":
    create_modern_template()
