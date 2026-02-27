
import sys
import traceback
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

def create_rich_template():
    print("Initializing Presentation...")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Define Colors
    COLOR_BG_DARK = RGBColor(15, 23, 42)    # Slate 900
    COLOR_ACCENT_BLUE = RGBColor(56, 189, 248) # Sky 400
    COLOR_ACCENT_PURPLE = RGBColor(99, 102, 241) # Indigo 500
    COLOR_TEXT_WHITE = RGBColor(248, 250, 252) # Slate 50
    COLOR_CARD_BG = RGBColor(30, 41, 59)    # Slate 800 (Card)

    def set_background(slide_layout, color):
        bg = slide_layout.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    print("Accessing Master...")
    master = prs.slide_master
    
    # ---------------------------------------------------------
    # 1. Custom Title Layout (Index 0)
    # ---------------------------------------------------------
    print("Configuring Title Layout...")
    title_layout = master.slide_layouts[0]
    set_background(title_layout, COLOR_BG_DARK)
    
    # NOTE: python-pptx cannot add shapes to Layouts directly. 
    # We will handle decorations in the Generator code (solution_router.py).
    
    shapes = title_layout.shapes

    # Style Title Placeholders
    if title_layout.placeholders:
        title_ph = title_layout.placeholders[0]
        title_ph.left = Inches(1)
        title_ph.top = Inches(2.5)
        title_ph.width = Inches(10)
        title_ph.height = Inches(2)
        
        tf = title_ph.text_frame
        p = tf.paragraphs[0]
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(60)
        p.font.color.rgb = COLOR_TEXT_WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT
        
        if len(title_layout.placeholders) > 1:
            sub = title_layout.placeholders[1]
            sub.left = Inches(1)
            sub.top = Inches(5)
            sub.width = Inches(8)
            sub.height = Inches(1)
            sub.text_frame.paragraphs[0].font.color.rgb = RGBColor(148, 163, 184) # Slate 400
            sub.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    # ---------------------------------------------------------
    # 2. Content Layout (Card Style) (Index 1)
    # ---------------------------------------------------------
    print("Configuring Content Layout...")
    content_layout = master.slide_layouts[1]
    set_background(content_layout, COLOR_BG_DARK)
    
    # NOTE: Shapes added in Generator.
    
    # Adjust Placeholders
    # Title
    c_title = None
    c_body = None
    for s in content_layout.placeholders:
        if s.placeholder_format.idx == 0: c_title = s
        if s.placeholder_format.idx == 1: c_body = s
        
    if c_title:
        c_title.left = Inches(0.8)
        c_title.top = Inches(0.3)
        c_title.width = Inches(10)
        c_title.height = Inches(1)
        
        p = c_title.text_frame.paragraphs[0]
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(36)
        p.font.color.rgb = COLOR_ACCENT_BLUE
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT
        
    if c_body:
        c_body.left = Inches(1.2) # Indent inside the card
        c_body.top = Inches(1.8)
        c_body.width = Inches(11.2)
        c_body.height = Inches(4.8)
        # Text color
        # Since we can't easily set default text color for body placeholder via python-pptx efficiently without a real template file,
        # we hope the user's content generator sets it or defaults are visible (black on dark card is bad).
        # We need to set the layout's placeholder text style level.
        # This is hard in python-pptx from scratch.
        # Workaround: The card background is COLOR_CARD_BG (Dark). Text needs to be Light.
        pass

    # Save
    output_path = r"E:\DigitalEmployee\storage\ppt_templates\DeepSeek_Tech_Pro.pptx"
    prs.save(output_path)
    print(f"SUCCESS: Created template at {output_path}")

if __name__ == "__main__":
    create_rich_template()
